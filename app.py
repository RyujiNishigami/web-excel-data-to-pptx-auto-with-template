import os
import uuid
import re
import copy
import io
import traceback
import zipfile
from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd

app = Flask(__name__)

def copy_run_format(src_run, dest_run):
    dest_run.font.bold = src_run.font.bold
    dest_run.font.italic = src_run.font.italic
    dest_run.font.underline = src_run.font.underline
    dest_run.font.size = src_run.font.size
    dest_run.font.name = src_run.font.name
    if src_run.font.color.type is not None:
        try:
            dest_run.font.color.rgb = src_run.font.color.rgb
        except:
            try:
                theme_color = src_run.font.color.theme_color
                if theme_color != 0:
                    dest_run.font.color.theme_color = theme_color
            except:
                pass

def safe_replace_paragraph(paragraph, new_text):
    if not paragraph.runs:
        paragraph.text = str(new_text)
        return
        
    first_run = paragraph.runs[0]
    
    font_name = first_run.font.name
    font_size = first_run.font.size
    font_bold = first_run.font.bold
    font_italic = first_run.font.italic
    font_underline = first_run.font.underline
    
    try: font_color_rgb = first_run.font.color.rgb
    except: font_color_rgb = None
    
    try: font_color_theme = first_run.font.color.theme_color
    except: font_color_theme = None
        
    for r in paragraph.runs:
        r.text = ""
            
    first_run.text = str(new_text)
    
    if font_name is not None: first_run.font.name = font_name
    if font_size is not None: first_run.font.size = font_size
    if font_bold is not None: first_run.font.bold = font_bold
    if font_italic is not None: first_run.font.italic = font_italic
    if font_underline is not None: first_run.font.underline = font_underline
    
    if font_color_rgb is not None:
        first_run.font.color.rgb = font_color_rgb
    elif font_color_theme is not None and font_color_theme != 0:
        first_run.font.color.theme_color = font_color_theme

def safe_replace_cell(cell, new_text):
    if not cell.text_frame.paragraphs:
        cell.text = str(new_text)
        return
    safe_replace_paragraph(cell.text_frame.paragraphs[0], str(new_text))
    for i in range(1, len(cell.text_frame.paragraphs)):
        p = cell.text_frame.paragraphs[i]
        for r in p.runs:
            r.text = ""

def duplicate_last_row(table):
    tbl = table._tbl
    last_tr = tbl.tr_lst[-1]
    new_tr = copy.deepcopy(last_tr)
    tbl.append(new_tr)

def delete_excess_rows(table, keep_count):
    tbl = table._tbl
    while len(tbl.tr_lst) > keep_count:
        tr = tbl.tr_lst[-1]
        tbl.remove(tr)

def update_ppt_table(table, tables_data, tag_idx):
    table_idx = tag_idx - 1
    if 0 <= table_idx < len(tables_data):
        df_table = tables_data[table_idx]
        if df_table.empty:
            raise ValueError(f"Table {{{{Table_{tag_idx}}}}} is empty in Excel data!")
            
        start_row = 1
        sal_text = " ".join([str(x) for x in df_table.columns])
        if len(df_table) > 0:
            sal_text += " " + " ".join([str(x) for x in df_table.iloc[0]])
            
        is_temp = '℃' in sal_text or 'C)' in sal_text or 'C ' in sal_text
        is_humid = '%RH' in sal_text or '40-59' in sal_text

        skip_col_indices = []
        for i, col in enumerate(df_table.columns):
            col_upper = str(col).strip().upper()
            if col_upper in ['SAL', 'TARGET', 'CRITERIA', 'STANDARD'] or 'UNNAMED' in col_upper:
                skip_col_indices.append(i)
                continue
            
            is_empty = True
            for val in df_table.iloc[:, i]:
                val_str = str(val).strip().lower()
                if val_str not in ['', 'nan', 'none', 'null']:
                    is_empty = False
                    break
            if is_empty:
                skip_col_indices.append(i)

        target_rows_needed = start_row + len(df_table)
        while len(table.rows) < target_rows_needed:
            duplicate_last_row(table)
        if len(table.rows) > target_rows_needed:
            delete_excess_rows(table, target_rows_needed)

        for data_r_idx, row_data in df_table.iterrows():
            target_row_idx = start_row + data_r_idx
            reference_row = table.rows[1] if len(table.rows) > 1 else table.rows[0]
            
            for data_c_idx, value in enumerate(row_data):
                if data_c_idx < len(table.columns):
                    if data_c_idx in skip_col_indices:
                        continue

                    target_cell = table.cell(target_row_idx, data_c_idx)
                    val_str = "" if pd.isna(value) else str(value).strip()
                    if val_str.lower() in ['nan', 'none', 'null']:
                        val_str = ""
                    
                    if data_c_idx > 1 and val_str != "":
                        try:
                            if is_humid:
                                val_float = float(val_str)
                                if val_float <= 1.0: val_str = f"{int(val_float * 100)}%"
                                else: val_str = f"{int(val_float)}%"
                            elif is_temp:
                                val_float = float(val_str)
                                if val_float.is_integer(): val_str = f"{int(val_float)} ℃"
                                else: val_str = f"{val_float} ℃"
                        except ValueError:
                            pass

                    safe_replace_cell(target_cell, val_str)
                    
                    if target_row_idx > 1 and target_cell.text_frame.paragraphs[0].runs:
                        if data_c_idx < len(reference_row.cells) and reference_row.cells[data_c_idx].text_frame.paragraphs and reference_row.cells[data_c_idx].text_frame.paragraphs[0].runs:
                            ref_run = reference_row.cells[data_c_idx].text_frame.paragraphs[0].runs[0]
                            target_run = target_cell.text_frame.paragraphs[0].runs[0]
                            copy_run_format(ref_run, target_run)
    else:
        raise IndexError(f"Table {{{{Table_{tag_idx}}}}} not found in uploaded Excel data! (Only {len(tables_data)} tables found)")

def update_ppt_chart(chart, tables_data, tag_idx):
    chart_idx = tag_idx - 1
    if 0 <= chart_idx < len(tables_data):
        df = tables_data[chart_idx].copy()
        if df.empty:
            raise ValueError(f"Chart {{{{Chart_{tag_idx}}}}} data is completely empty in Excel!")
            
        sal_text = " ".join([str(x) for x in df.columns])
        if len(df) > 0:
            sal_text += " " + " ".join([str(x) for x in df.iloc[0]])
            
        is_temp = '℃' in sal_text or 'C)' in sal_text or 'C ' in sal_text
        is_humid = '%RH' in sal_text or '40-59' in sal_text
        
        # ถอดระบบซ่อนเดือนทิ้ง (ไม่ต้องหา Empty Cols แล้ว) เพื่อให้เดือนที่ว่างแสดงบนแกนกราฟเสมอ
        cols_to_drop = []
        for col in df.columns:
            col_str = str(col).strip().upper()
            if col_str in ['SAL', 'TARGET', 'CRITERIA', 'STANDARD'] or 'UNNAMED' in col_str:
                cols_to_drop.append(col)
            # เอาคอลัมน์ที่ไม่ใช่เดือน (เช่น (18 ºC - 26 ºC)) ออกจากกราฟ
            elif '°C' in col_str or 'ºC' in col_str or ' C ' in col_str or 'C)' in col_str or '%RH' in col_str:
                cols_to_drop.append(col)
                
        if cols_to_drop:
            df = df.drop(columns=cols_to_drop)

        if len(df.columns) < 2:
            raise ValueError(f"Chart {{{{Chart_{tag_idx}}}}} doesn't have enough columns to plot data. Missing numeric values.")

        chart_data = CategoryChartData()
        categories = [str(c).strip() for c in df.columns[1:]]
        chart_data.categories = categories

        for r_idx, row in df.iterrows():
            series_name = str(row.iloc[0]).strip()
            series_values = []
            
            for val in row.iloc[1:]:
                if pd.isna(val):
                    series_values.append(None)
                    continue
                try:
                    raw_val = str(val).replace(',', '').strip()
                    if raw_val == '' or raw_val.lower() in ['nan', 'none', 'null']:
                        series_values.append(None)
                    else:
                        if '%' in raw_val:
                            val_float = float(raw_val.replace('%', '')) / 100.0
                        else:
                            val_float = float(raw_val)
                        series_values.append(val_float)
                except Exception as e:
                    series_values.append(None)
                    
            chart_data.add_series(series_name, tuple(series_values))
            
        try:
            chart.replace_data(chart_data)
        except Exception as e:
            raise RuntimeError(f"Chart {{{{Chart_{tag_idx}}}}}: PowerPoint rejected the data structure. Error: {e}")
        
        for series in chart.series:
            try:
                series.has_data_labels = True
                if is_humid:
                    series.data_labels.number_format = '0%'
                    series.data_labels.number_format_is_linked = False
                elif is_temp:
                    series.data_labels.number_format = '0" °C"' 
                    series.data_labels.number_format_is_linked = False
            except Exception as e:
                pass
    else:
        raise IndexError(f"Chart {{{{Chart_{tag_idx}}}}} not found in uploaded Excel data! (Only {len(tables_data)} tables found)")

def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape

def extract_vars_from_shapes(shapes, pptx_vars):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            extract_vars_from_shapes(shape.shapes, pptx_vars)
        else:
            if shape.name:
                matches = re.findall(r'\{\{([^}]+)\}\}', shape.name)
                for m in matches:
                    if not re.match(r'^(table|chart)[\s_]*\d+$', m.strip(), re.IGNORECASE):
                        pptx_vars.add('{{' + m.strip() + '}}')
            
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    matches = re.findall(r'\{\{([^}]+)\}\}', p.text)
                    for m in matches:
                        if not re.match(r'^(table|chart)[\s_]*\d+$', m.strip(), re.IGNORECASE):
                            pptx_vars.add('{{' + m.strip() + '}}')
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        matches = re.findall(r'\{\{([^}]+)\}\}', cell.text)
                        for m in matches:
                            if not re.match(r'^(table|chart)[\s_]*\d+$', m.strip(), re.IGNORECASE):
                                pptx_vars.add('{{' + m.strip() + '}}')
            if shape.has_chart and shape.chart.has_title and shape.chart.chart_title.text_frame:
                matches = re.findall(r'\{\{([^}]+)\}\}', shape.chart.chart_title.text_frame.text)
                for m in matches:
                    if not re.match(r'^(table|chart)[\s_]*\d+$', m.strip(), re.IGNORECASE):
                        pptx_vars.add('{{' + m.strip() + '}}')

def replace_smart(prs, variables_dict, tables_data):
    for slide in prs.slides:
        table_updates = []
        chart_updates = []
        
        all_shapes = list(iter_shapes(slide.shapes))
        
        for shape in all_shapes:
            table_tag_idx = None
            chart_tag_idx = None
            
            if shape.name:
                m_name = re.findall(r'\{\{\s*(table|chart)[\s_]*(\d+)\s*\}\}', shape.name, re.IGNORECASE)
                for m in m_name:
                    if m[0].lower() == 'table': table_tag_idx = int(m[1])
                    if m[0].lower() == 'chart': chart_tag_idx = int(m[1])

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    p_text = paragraph.text
                    changed = False
                    for key, value in variables_dict.items():
                        if key in p_text:
                            p_text = p_text.replace(key, str(value))
                            changed = True
                            
                    matches = re.findall(r'\{\{\s*(table|chart)[\s_]*(\d+)\s*\}\}', p_text, re.IGNORECASE)
                    for m in matches:
                        if m[0].lower() == 'table': table_tag_idx = int(m[1])
                        if m[0].lower() == 'chart': chart_tag_idx = int(m[1])
                        changed = True
                        
                    if changed:
                        clean_text = re.sub(r'\{\{\s*(table|chart)[\s_]*(\d+)\s*\}\}', '', p_text, flags=re.IGNORECASE).strip()
                        safe_replace_paragraph(paragraph, clean_text)
                            
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        p_text = cell.text
                        changed = False
                        for key, value in variables_dict.items():
                            if key in p_text:
                                p_text = p_text.replace(key, str(value))
                                changed = True
                                
                        matches = re.findall(r'\{\{\s*(table|chart)[\s_]*(\d+)\s*\}\}', p_text, re.IGNORECASE)
                        for m in matches:
                            if m[0].lower() == 'table': table_tag_idx = int(m[1])
                            if m[0].lower() == 'chart': chart_tag_idx = int(m[1])
                            changed = True
                        
                        if changed:
                            clean_text = re.sub(r'\{\{\s*(table|chart)[\s_]*(\d+)\s*\}\}', '', p_text, flags=re.IGNORECASE).strip()
                            safe_replace_cell(cell, clean_text)
                            
            if shape.has_chart and shape.chart.has_title and shape.chart.chart_title.text_frame:
                title_text = shape.chart.chart_title.text_frame.text
                changed = False
                for key, value in variables_dict.items():
                    if key in title_text:
                        title_text = title_text.replace(key, str(value))
                        changed = True
                
                matches = re.findall(r'\{\{\s*(table|chart)[\s_]*(\d+)\s*\}\}', title_text, re.IGNORECASE)
                for m in matches:
                    if m[0].lower() == 'table': table_tag_idx = int(m[1])
                    if m[0].lower() == 'chart': chart_tag_idx = int(m[1])
                    changed = True
                    
                if changed:
                    clean_text = re.sub(r'\{\{\s*(table|chart)[\s_]*(\d+)\s*\}\}', '', title_text, flags=re.IGNORECASE).strip()
                    if shape.chart.chart_title.text_frame.paragraphs:
                        safe_replace_paragraph(shape.chart.chart_title.text_frame.paragraphs[0], clean_text)
                    else:
                        shape.chart.chart_title.text_frame.text = clean_text

            if table_tag_idx is not None and shape.has_table:
                shape.name = f"{{{{Table_{table_tag_idx}}}}}"
                table_updates.append((shape, table_tag_idx))
                
            if chart_tag_idx is not None and shape.has_chart:
                shape.name = f"{{{{Chart_{chart_tag_idx}}}}}"
                chart_updates.append((shape, chart_tag_idx))

        for shape, tag_idx in table_updates:
            try:
                update_ppt_table(shape.table, tables_data, tag_idx)
            except Exception as e:
                print(f"Skipped Table {tag_idx} update due to error: {e}")
                
        for shape, tag_idx in chart_updates:
            try:
                update_ppt_chart(shape.chart, tables_data, tag_idx)
            except Exception as e:
                print(f"Skipped Chart {tag_idx} update due to error: {e}")

UPLOADED_PPTX_MEMORY = {}

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/preview', methods=['POST'])
def preview_data():
    excel_file = request.files.get('excel_file')
    pptx_files = request.files.getlist('pptx_file')
    
    if not excel_file or excel_file.filename == '' or not pptx_files or pptx_files[0].filename == '':
        return "Please upload Excel and at least one PPTX file.", 400
        
    excel_bytes = excel_file.read()
    
    memory_id = uuid.uuid4().hex
    
    saved_pptx = []
    for pf in pptx_files:
        if pf.filename:
            saved_pptx.append({
                'filename': pf.filename,
                'bytes': pf.read()
            })
            
    UPLOADED_PPTX_MEMORY[memory_id] = saved_pptx
    
    variables_dict = {}
    tables = []

    try:
        df = pd.read_excel(io.BytesIO(excel_bytes), sheet_name=0, header=None)
        
        for r in range(len(df)):
            for c in range(len(df.columns)):
                val = str(df.iloc[r, c]).strip()
                if val.startswith('{{') and val.endswith('}}'):
                    if c + 1 < len(df.columns):
                        next_val = df.iloc[r, c+1]
                        variables_dict[val] = str(next_val) if pd.notna(next_val) else ""

        current_table_rows = []
        start_col, end_col = -1, -1
        
        for r in range(len(df)):
            row = df.iloc[r]
            non_null_cols = [c for c, val in enumerate(row) if pd.notna(val) and str(val).strip() != ""]
            
            if len(non_null_cols) >= 2:
                first_val = str(row[non_null_cols[0]]).strip()
                is_var = first_val.startswith('{{') and first_val.endswith('}}')
                
                if not is_var:
                    if not current_table_rows:
                        start_col = min(non_null_cols)
                        end_col = max(non_null_cols)
                    
                    cleaned_row = [str(x).strip() if pd.notna(x) else "" for x in row[start_col:end_col+1]]
                    current_table_rows.append(cleaned_row)
                else:
                    if current_table_rows:
                        tables.append(current_table_rows)
                        current_table_rows = []
            else:
                if current_table_rows:
                    tables.append(current_table_rows)
                    current_table_rows = []
                    
        if current_table_rows:
            tables.append(current_table_rows)
            
    except Exception as e:
        print(f"Error parsing excel: {e}")

    try:
        pptx_vars = set()
        for ppt_data in saved_pptx:
            prs = Presentation(io.BytesIO(ppt_data['bytes']))
            for slide in prs.slides:
                extract_vars_from_shapes(slide.shapes, pptx_vars)
            
        for v in pptx_vars:
            if v not in variables_dict:
                variables_dict[v] = "" 
    except Exception as e:
        print(f"Error parsing PPTX for variables: {e}")

    return render_template('edit.html', variables=variables_dict, tables=tables, pptx_filename=memory_id)

@app.route('/generate', methods=['POST'])
def generate_report():
    memory_id = request.form.get('pptx_filename')
    
    if memory_id not in UPLOADED_PPTX_MEMORY:
        return "Error: PowerPoint templates not found in memory (Session expired). Please re-upload.", 400
        
    saved_pptx = UPLOADED_PPTX_MEMORY[memory_id]
    
    variables_dict = {}
    for key, value in request.form.items():
        if key.startswith('var_'):
            real_key = key.replace('var_', '')
            variables_dict[real_key] = value
            
    tables = []
    table_count = int(request.form.get('table_count', 0))
    for t in range(table_count):
        row_count = int(request.form.get(f't{t}_rows', 0))
        col_count = int(request.form.get(f't{t}_cols', 0))
        
        table_data = []
        headers = []
        for r in range(row_count):
            row_data = []
            for c in range(col_count):
                cell_val = request.form.get(f'table_{t}_{r}_{c}', '')
                row_data.append(cell_val)
            if r == 0:
                headers = row_data
            else:
                table_data.append(row_data)
        
        if headers or table_data:
            df_t = pd.DataFrame(table_data, columns=headers if headers else None)
            tables.append(df_t)

    try:
        if len(saved_pptx) == 1:
            ppt_data = saved_pptx[0]
            prs = Presentation(io.BytesIO(ppt_data['bytes']))
            replace_smart(prs, variables_dict, tables)
            
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            del UPLOADED_PPTX_MEMORY[memory_id]
            return send_file(output, as_attachment=True, download_name=ppt_data['filename'])
            
        else:
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for ppt_data in saved_pptx:
                    prs = Presentation(io.BytesIO(ppt_data['bytes']))
                    replace_smart(prs, variables_dict, tables)
                    
                    output = io.BytesIO()
                    prs.save(output)
                    
                    zipf.writestr(ppt_data['filename'], output.getvalue())
            
            zip_buffer.seek(0)
            del UPLOADED_PPTX_MEMORY[memory_id]
            return send_file(zip_buffer, as_attachment=True, download_name="Automated_Reports.zip", mimetype="application/zip")
            
    except Exception as e:
        error_trace = traceback.format_exc()
        print(error_trace)
        error_html = (
            "<div style='font-family: Arial, sans-serif; padding: 20px;'>"
            "<h3 style='color: #d9534f;'>🚨 ระบบแจ้งเตือนจุดพัง (Diagnostic Mode)</h3>"
            "<p>พบข้อผิดพลาดในขณะแก้ไขไฟล์ PPTX:</p>"
            f"<pre style='background: #f4f4f4; padding: 15px; border-radius: 5px; font-size: 14px; color: #333; overflow-x: auto; border: 1px solid #ddd; white-space: pre-wrap;'>{error_trace}</pre>"
            "<br>"
            "<a href='/' style='padding: 10px 15px; background: #0275d8; color: white; text-decoration: none; border-radius: 4px;'>กลับไปหน้าแรก</a>"
            "</div>"
        )
        return error_html, 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)